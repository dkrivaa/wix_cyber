import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_TICK_MARK
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
import datetime
import numpy as np


# Extracting annual data from dataframe for slides
def annual(df):
    # Adding year column to df
    df['year'] = np.ceil(df['month'] / 12).astype(int)
    customers = df.groupby('year')['customers'].sum()
    risk_packages = df.groupby('year')['Risk assessment pakages sold'].sum()
    soc_packages = df.groupby('year')['SOC pakages sold'].sum()
    insurance_packages = df.groupby('year')['Insurance packages sold'].sum()
    income = df.groupby('year')['Total income'].sum()
    cost = df.groupby('year')['Total cost'].sum()
    profit = df.groupby('year')['Gross profit'].sum()

    return customers, risk_packages, soc_packages, insurance_packages, income, cost, profit


def make_slides(df, data_dict, file_name):
    # create presentation
    prs = Presentation()

    ### SLIDE FRONT ##############################################

    front = prs.slides.add_slide(prs.slide_layouts[6])
    # TEXTBOX
    textbox = front.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    # Add text to the textbox
    p1 = text_frame.add_paragraph()
    p1.text = f"CyberMarket Simulation: {data_dict['scenarioName']}"
    p1.font.size = Pt(32)
    p1.font.color.rgb = RGBColor(16, 53, 117)
    p1.font.bold = True
    p1.alignment = 2

    p2 = text_frame.add_paragraph()
    my_time = datetime.datetime.now() + datetime.timedelta(hours=3)
    p2.text = f"{my_time.strftime("%Y-%m-%d %H:%M:%S")}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(16, 53, 117)
    p2.font.bold = True
    p2.alignment = 2



    ### SLIDE 1 ##############################################
    # Scenario assumptions

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # TEXTBOX
    textbox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.5), Inches(3))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    # Add text to the textbox
    p1 = text_frame.add_paragraph()
    p1.text = "Main Scenario Assumptions:"
    p1.font.size = Pt(28)
    p1.font.color.rgb = RGBColor(16, 53, 117)
    p1.font.bold = True
    p1.alignment = 1
    # Policy price
    p2 = text_frame.add_paragraph()
    p2.text = f'Policy Price: ₪{data_dict['insurancePrice']:,.0f}'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(16, 53, 117)
    p2.alignment = 1
    # Commission Rates
    p3 = text_frame.add_paragraph()
    p3.text = (f'Commission Rates: New: {data_dict['newCommission']:.0f}%, Referred: {data_dict['referredCommission']:.0f}%, '
               f'Lead: {data_dict['leadCommission']:.0f}%, Existing: {data_dict['existingCommission']:.0f}%')
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(16, 53, 117)
    p3.alignment = 1

    # Salaries
    p4 = text_frame.add_paragraph()
    p4.text = (f'Salaries: Admin: ₪{data_dict['adminSalary']:,.0f}, Tele meeting: ₪{data_dict['teleSalary']:,.0f}, '
               f'Sales: ₪{data_dict['salesSalary']:,.0f}, Analyst: ₪{data_dict['cyberSalary']:,.0f}, '
               f'Logistics: ₪{data_dict['logisticsSalary']:,.0f}')
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(16, 53, 117)
    p4.alignment = 1
    # Incentives
    p5 = text_frame.add_paragraph()
    p5.text = (f'Bonuses: Schedule Sales Meeting: ₪{data_dict['teleIncentive']:,.0f}, '
               f'Policy Sale: {data_dict['salesIncentive']:.0f}% of commission')
    p5.font.size = Pt(16)
    p5.font.color.rgb = RGBColor(16, 53, 117)
    p5.alignment = 1
    # Cost of Lead (buying Policy)
    p6 = text_frame.add_paragraph()
    p6.text = (f' Cost of lead: ₪{data_dict['leadCost']}')
    p6.font.size = Pt(18)
    p6.font.color.rgb = RGBColor(16, 53, 117)
    p6.alignment = 1
    # Perceived risk
    p7 = text_frame.add_paragraph()
    p7.text = (
        f'Risk Perception: New: {data_dict['newRisk']:.0f}%, Referred: {data_dict['referredRisk']:.0f}%, '
        f'Lead: {data_dict['leadRisk']:.0f}%, Existing: {data_dict['existingRisk']:.0f}%')
    p7.font.size = Pt(18)
    p7.font.color.rgb = RGBColor(16, 53, 117)
    p7.alignment = 1

    # Set the line spacing
    for p in text_frame.paragraphs:
        p.line_spacing = 2

    ### ANNUAL SLIDES ########################################

    # Annual Charts
    # Making annual data series for slides 4-6
    customers, risk_packages, soc_packages, insurance_packages, income, cost, profit = annual(df)
    # Making year list
    years = []
    for i in range(len(customers)):
        year = f'year{i + 1}'
        years.append(year)

    ### SLIDE 2 ##############################################

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - Gross Profit

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('Gross Profit', profit)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide2.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Annual Gross Profit"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 3 ##############################################

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - Insurance Packages Sold

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('insurance packages', insurance_packages)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide3.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Insurance Packages Sold"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 4 ##############################################

    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - Risk Assessment Packages Sold

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('Risk packages', risk_packages)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide4.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Risk Assessment Packages Sold"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False


    ### SLIDE 5 ##############################################

    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - SOC Packages Sold

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('SOC packages', soc_packages)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "SOC Packages Sold"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False


    ### SLIDE 6 ##############################################

    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - Annual Income

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('income', income)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide6.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Annual Income"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False


    ### SLIDE 7 ##############################################

    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - Annual Expenses

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('cost', cost)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide7.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Annual Expenses"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False


    ### SLIDE 8 ##############################################

    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - Annual Customers

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('customers', customers)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide8.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Customers"
    # Add legend
    chart.has_legend = False

    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 39 - i * r))
        green = max(0, min(255, 63 - i * g))
        blue = max(0, min(255, 92 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 9 ##############################################
    # Charts

    slide9 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - CUSTOMERS

    # Data for chart 1
    monthly_profit = df['Gross profit']

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = list(df['month'])
    chart_data.add_series('Gross Profit', monthly_profit)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide9.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Monthly Gross Profit"
    # Add legend
    chart.has_legend = False
    # Position legend at bottom
    chart.legend.position = XL_LEGEND_POSITION.TOP
    # Set font size for the legend
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = True
    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 39
        g = 63
        b = 92

        red = max(0, min(255, 151 - i * r))
        green = max(0, min(255, 228 - i * g))
        blue = max(0, min(255, 100 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False


    ### SLIDE 10 ##############################################
    # Charts

    slide10 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - CUSTOMERS

    # Data for chart 1
    acc_profit = df['Acc. Gross Profit']

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = list(df['month'])
    chart_data.add_series('Acc. Gross Profit', acc_profit)

    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide10.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Accumulated Gross Profit"
    # Add legend
    chart.has_legend = False
    # Position legend at bottom
    chart.legend.position = XL_LEGEND_POSITION.TOP
    # Set font size for the legend
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = True
    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 39
        g = 63
        b = 92

        red = max(0, min(255, 151 - i * r))
        green = max(0, min(255, 228 - i * g))
        blue = max(0, min(255, 100 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 11 ##############################################
    # Charts

    slide11 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - CUSTOMERS

    # Data for chart 1
    new_customer = df['new customers']
    referred_customer = df['referred customers']
    lead_customer = df['lead customers']
    existing_customer = df['existing customers']

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = list(df['month'])
    chart_data.add_series('new', new_customer)
    chart_data.add_series('referred', referred_customer)
    chart_data.add_series('lead', lead_customer)
    chart_data.add_series('existing', existing_customer)
    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide11.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Customers By Type"
    # Add legend
    chart.has_legend = True
    # Position legend at bottom
    chart.legend.position = XL_LEGEND_POSITION.TOP
    # Set font size for the legend
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = True
    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 151 - i * r))
        green = max(0, min(255, 228 - i * g))
        blue = max(0, min(255, 100 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 12 ##############################################
    # Charts

    slide12 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - CUSTOMERS buying insurance package

    # Data for chart
    new_customer = df['Insurance packages sold to new customers']
    referred_customer = df['Insurance packages sold to referred customers']
    lead_customer = df['Insurance packages sold to lead customers']
    existing_customer = df['Insurance packages sold to existing customers']

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = list(df['month'])
    chart_data.add_series('new', new_customer)
    chart_data.add_series('referred', referred_customer)
    chart_data.add_series('lead', lead_customer)
    chart_data.add_series('existing', existing_customer)
    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide12.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Customers Buying Insurance Package By Type"
    # Add legend
    chart.has_legend = True
    # Position legend at bottom
    chart.legend.position = XL_LEGEND_POSITION.TOP
    # Set font size for the legend
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = True
    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 151 - i * r))
        green = max(0, min(255, 228 - i * g))
        blue = max(0, min(255, 100 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 13 ##############################################
    # Charts

    slide13 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - CUSTOMERS buying insurance package

    # Data for chart
    new_customer = df['Risk assessment pakages sold to new customers']
    referred_customer = df['Risk assessment pakages sold to referred customers']
    lead_customer = df['Risk assessment pakages sold to lead customers']
    existing_customer = df['Risk assessment pakages sold to existing customers']

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = list(df['month'])
    chart_data.add_series('new', new_customer)
    chart_data.add_series('referred', referred_customer)
    chart_data.add_series('lead', lead_customer)
    chart_data.add_series('existing', existing_customer)
    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide13.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Customers Buying Risk Assessment Package By Type"
    # Add legend
    chart.has_legend = True
    # Position legend at bottom
    chart.legend.position = XL_LEGEND_POSITION.TOP
    # Set font size for the legend
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = True
    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 151 - i * r))
        green = max(0, min(255, 228 - i * g))
        blue = max(0, min(255, 100 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False

    ### SLIDE 14 ##############################################
    # Charts

    slide14 = prs.slides.add_slide(prs.slide_layouts[6])

    # CHART - CUSTOMERS buying insurance package

    # Data for chart
    new_customer = df['SOC pakages sold to new customers']
    referred_customer = df['SOC pakages sold to referred customers']
    lead_customer = df['SOC pakages sold to lead customers']
    existing_customer = df['SOC pakages sold to existing customers']

    # define chart and data
    chart_data = CategoryChartData()
    chart_data.categories = list(df['month'])
    chart_data.add_series('new', new_customer)
    chart_data.add_series('referred', referred_customer)
    chart_data.add_series('lead', lead_customer)
    chart_data.add_series('existing', existing_customer)
    # add chart to slide
    x, y, chart_x, chart_y = Inches(0.5), Inches(1), Inches(9), Inches(6)
    chart_frame = slide14.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, chart_x, chart_y, chart_data)

    # Define chart object to manipulate the object characteristics
    chart = chart_frame.chart
    # Add title to chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Customers Buying SOC Package By Type"
    # Add legend
    chart.has_legend = True
    # Position legend at bottom
    chart.legend.position = XL_LEGEND_POSITION.TOP
    # Set font size for the legend
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = True
    # Customize colors of columns
    for i in range(len(chart.series)):
        series = chart.series[i]
        fill = series.format.fill
        fill.solid()
        r = 38
        g = 58
        b = 26

        red = max(0, min(255, 151 - i * r))
        green = max(0, min(255, 228 - i * g))
        blue = max(0, min(255, 100 - i * b))
        fill.fore_color.rgb = RGBColor(red, green, blue)

    # Set font size for category (X) axis
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)  # Set font size to 8 points
    # Set font size for value (Y) axis
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)  # Set font size to 14 points
    # Set number format
    value_axis.tick_labels.number_format = '#,##0'
    # Remove gridlines
    value_axis.has_major_gridlines = False


    # Save the presentation to the given file path
    prs.save(file_name)